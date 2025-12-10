import os
import requests
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


BG_COLOR = RGBColor(30, 30, 30)      
TITLE_COLOR = RGBColor(255, 215, 0)  
TEXT_COLOR = RGBColor(230, 230, 230)  
ACCENT_COLOR = RGBColor(70, 130, 180)

def download_image(url, save_path):
    if not url or not isinstance(url, str) or not url.startswith("http"):
        return None
    try:
        r = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        r.raise_for_status()
        with open(save_path, "wb") as f:
            f.write(r.content)
        return save_path
    except:
        return None

def set_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def split_text_into_bullets(raw_content):
    """
    Splits string into bullets based on COMMAS, PERIODS, or NEWLINES.
    """
    if isinstance(raw_content, list):
        return raw_content
    
    if not isinstance(raw_content, str):
        return []

    text = raw_content.replace('\n', ',')

    parts = re.split(r'[.,]\s*', text)
    
    final_bullets = [p.strip() for p in parts if p.strip()]
                
    return final_bullets

def create_ppt(slides_data, output_path="presentation.pptx", topic="AI Presentation"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.text = topic
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    for i, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = slide_info.get("title", "Slide Title")
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = TITLE_COLOR

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.fill.background()

        image_url = slide_info.get("image_url")
        has_image = bool(image_url)
        
        if has_image:
            TEXT_WIDTH = Inches(7.5)
        else:
            TEXT_WIDTH = Inches(12.33)

        raw_content = slide_info.get("content", [])
        bullet_points = split_text_into_bullets(raw_content)

        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), TEXT_WIDTH, Inches(5.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.clear() 

        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point 
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(14) 
            
            p.level = 0 

        if has_image:
            temp_img = f"temp_{i}.jpg"
            img_path = download_image(image_url, temp_img)
            if img_path:
                try:
                    pic = slide.shapes.add_picture(img_path, Inches(8.5), Inches(1.8))
                    max_w, max_h = Inches(4.3), Inches(5.0)
                    ratio = pic.width / pic.height
                    if ratio > (max_w / max_h):
                        pic.width = max_w
                        pic.height = int(max_w / ratio)
                    else:
                        pic.height = max_h
                        pic.width = int(max_h * ratio)
                    
                    pic.left = int(Inches(8.5) + ((max_w - pic.width) / 2))
                    pic.top = int(Inches(1.8) + ((max_h - pic.height) / 2))
                except:
                    pass
                if os.path.exists(img_path): os.remove(img_path)

    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

