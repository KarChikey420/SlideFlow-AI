from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import requests
import json

def create_ppt_from_result(result, output_filename="presentation.pptx"):
    # Try to get slides from the 'slides' key first (new format)
    if 'slides' in result:
        slides_data = result['slides']
    else:
        # Fallback to extracting from messages
        slides_data = extract_slides_from_result(result)
    
    if not slides_data:
        print("❌ No slides data found!")
        return
    
    print(f"📊 Creating presentation with {len(slides_data)} slides...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for i, slide_data in enumerate(slides_data, 1):
        print(f"  📄 Adding slide {i}: {slide_data.get('title', 'Untitled')}")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', 'Untitled')
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Add bullets
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for bullet in slide_data.get('bullets', []):
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
        
        # Add image
        if 'image_url' in slide_data and slide_data['image_url']:
            try:
                print(f"    🖼️  Adding image...")
                response = requests.get(slide_data['image_url'], timeout=10)
                response.raise_for_status()
                img_stream = BytesIO(response.content)
                slide.shapes.add_picture(
                    img_stream, 
                    Inches(5.5), Inches(1.5), 
                    width=Inches(4), height=Inches(5)
                )
                print(f"    ✅ Image added")
            except Exception as e:
                print(f"    ⚠️  Could not add image: {e}")
        else:
            print(f"    ⚠️  No image URL found for this slide")
    
    prs.save(output_filename)
    print(f"\n✅ Presentation saved as {output_filename}")


def extract_slides_from_result(result):
    """
    Extract slides list from agent result.
    """
    try:
        if 'messages' not in result:
            return []
        
        # Find all AI messages
        from langchain_core.messages import AIMessage
        ai_messages = [msg for msg in result['messages'] if isinstance(msg, AIMessage)]
        
        if not ai_messages:
            return []
        
        # Try to find slides in AI messages (reverse order)
        for msg in reversed(ai_messages):
            try:
                content = msg.content if hasattr(msg, 'content') else str(msg)
                content = content.strip()
                
                # Remove markdown code blocks
                if content.startswith('```'):
                    lines = content.split('\n')
                    content = '\n'.join(lines[1:-1]) if len(lines) > 2 else content
                
                slides = json.loads(content)
                
                # Validate it's slide data
                if isinstance(slides, list) and len(slides) > 0:
                    if 'title' in slides[0]:
                        return slides
                elif isinstance(slides, dict) and 'title' in slides:
                    return [slides]
            except:
                continue
        
        return []
            
    except Exception as e:
        print(f"❌ Error extracting slides: {e}")
        return []


if __name__ == "__main__":
    # Import generator function
    from generator import generate_full_presentation
    
    # Generate presentation
    print("🚀 Generating presentation...\n")
    result = generate_full_presentation("AI in Healthcare", slide_count=5)
    
    print("\n" + "="*50)
    print("DEBUG INFO:")
    print("="*50)
    print("Result keys:", result.keys())
    print("\nNumber of messages:", len(result['messages']))
    
    # Show all message types
    print("\nMessage types:")
    for i, msg in enumerate(result['messages']):
        print(f"  {i}: {type(msg).__name__}")
    
    if 'slides' in result:
        print(f"\n✅ Found {len(result['slides'])} slides in result['slides']")
    print("="*50 + "\n")
    
    # Create PowerPoint
    create_ppt_from_result(result, "ai_healthcare_presentation.pptx")